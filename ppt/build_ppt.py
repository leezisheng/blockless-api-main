# -*- coding: utf-8 -*-
"""Blockless deck -> 10-page PPTX. White/gray text on a faded dark background."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.join(os.path.dirname(ROOT), "docs")
THUMB = os.path.join(ROOT, "_thumbs")
BG_169 = os.path.join(ROOT, "_bg16x9.png")
OUT = os.path.join(ROOT, "Blockless.pptx")
MOVIE = r"G:\公司相关\融资相关\3D仿真Demo视频\3D仿真演示.mp4"


def D(name):
    return os.path.join(DOCS, name)


# ---- palette : monochrome white/grey, brightened ------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY  = RGBColor(0xED, 0xF2, 0xFA)   # near-white body
GREY  = RGBColor(0xCE, 0xD7, 0xE4)   # secondary
GREY2 = RGBColor(0xA6, 0xB4, 0xC8)   # labels / captions / dim
LINE  = RGBColor(0x4A, 0x5B, 0x74)   # subtle borders
BG_DARK = RGBColor(0x05, 0x0A, 0x14)
BG_ALPHA = 55000                     # bg image opacity 55%

CJK = "Microsoft YaHei"
MONO = "Consolas"
DISP = "Segoe UI Semibold"

SWIN, SHIN = 13.333, 7.5


def IN(x):
    return Inches(x)


# ---- image prep ---------------------------------------------------------
def make_bg():
    im = Image.open(D("背景.png")).convert("RGB")
    w, h = im.size
    tw = int(round(h * 16 / 9))
    if tw <= w:
        im = im.crop((0, 0, tw, h))           # left-anchored, keep glow
    else:
        im = im.crop((0, 0, w, int(round(w * 9 / 16))))
    im.resize((1920, 1080), Image.LANCZOS).save(BG_169)


def cover_crop(src, dst, tw, th):
    im = Image.open(src).convert("RGB")
    w, h = im.size
    tr, r = tw / th, w / h
    if r > tr:
        nw = int(h * tr); x = (w - nw) // 2; im = im.crop((x, 0, x + nw, h))
    else:
        nh = int(w / tr); y = (h - nh) // 2; im = im.crop((0, y, w, y + nh))
    im.resize((tw, th), Image.LANCZOS).save(dst)


def make_thumbs():
    if not os.path.isdir(THUMB):
        os.makedirs(THUMB)
    for i in range(1, 10):
        cover_crop(D("用户案例%d.jpg" % i),
                   os.path.join(THUMB, "case%d.jpg" % i), 640, 480)


# ---- low level ----------------------------------------------------------
def set_cjk(run, name=CJK):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)


def set_spc(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def shadow_off(shape):
    spPr = shape._element.spPr
    if spPr.find(qn("a:effectLst")) is None:
        spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


def fill_alpha(shape, pct):
    sf = shape._element.spPr.find(qn("a:solidFill"))
    clr = sf.find(qn("a:srgbClr"))
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def panel(shape, rgb=WHITE, alpha_pct=5, line=LINE, line_w=1.0):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb
    fill_alpha(shape, alpha_pct)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(line_w)
    shadow_off(shape)


def set_pic_alpha(pic, amt=BG_ALPHA):
    blip = pic._element.find(qn("p:blipFill")).find(qn("a:blip"))
    blip.append(blip.makeelement(qn("a:alphaModFix"), {"amt": str(amt)}))


def run_style(run, size, color, bold=True, latin=DISP, spc=None, cjk=CJK):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = latin
    set_cjk(run, cjk)
    if spc is not None:
        set_spc(run, spc)


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def put(tf, runs, first=False, align=PP_ALIGN.CENTER, line=1.15, sb=0, sa=0):
    """runs: list of (text, size, color, bold, latin[, cjk])"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after = Pt(sa)
    try: p.line_spacing = line
    except Exception: pass
    for spec in runs:
        r = p.add_run(); r.text = spec[0]
        cjk = spec[5] if len(spec) > 5 else CJK
        run_style(r, spec[1], spec[2], spec[3], spec[4], cjk=cjk)
    return p


# ---- building blocks ----------------------------------------------------
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG_DARK
    set_pic_alpha(s.shapes.add_picture(BG_169, 0, 0, IN(SWIN), IN(SHIN)))
    return s


def add_label(s, text, y=0.52):
    tf = textbox(s, 0, y, SWIN, 0.32)
    put(tf, [("// " + text, 12.5, GREY2, False, MONO)], first=True)
    p = tf.paragraphs[0]
    for r in p.runs:
        set_spc(r, 4)


def add_title(s, runs, y=1.0, size=32, w=SWIN, l=0, line=1.12):
    tf = textbox(s, l, y, w, 1.0)
    specs = [(t, size, c, True, DISP) for (t, c) in runs]
    put(tf, specs, first=True, line=line)


def add_sep(s, y=1.86, w=0.9):
    sep = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             IN(SWIN / 2 - w / 2), IN(y), IN(w), IN(0.05))
    sep.fill.solid(); sep.fill.fore_color.rgb = GREY
    sep.line.fill.background(); shadow_off(sep)


def add_lead(s, runs, y, size=17, w=10.5, h=0.9, line=1.5):
    tf = textbox(s, SWIN / 2 - w / 2, y, w, h)
    specs = [(t, size, c, False, CJK) for (t, c) in runs]
    put(tf, specs, first=True, line=line)


def add_card(s, l, t, w, h, head, body, head_color=WHITE,
             head_size=16, body_size=12.5, corner=0.09):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    try: sh.adjustments[0] = corner
    except Exception: pass
    panel(sh)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = IN(0.22)
    tf.margin_top = IN(0.18); tf.margin_bottom = IN(0.14)
    put(tf, [(head, head_size, head_color, True, DISP)], first=True,
        align=PP_ALIGN.LEFT, line=1.05)
    put(tf, [(body, body_size, GREY, False, CJK)], align=PP_ALIGN.LEFT,
        line=1.32, sb=6)
    return sh


def add_image(s, box_l, box_t, box_w, box_h, path, caption):
    im = Image.open(path); ar = im.size[0] / im.size[1]
    cap_h = 0.30
    avail = box_h - cap_h
    if box_w / avail > ar:
        h = avail; w = h * ar
    else:
        w = box_w; h = w / ar
    l = box_l + (box_w - w) / 2
    t = box_t + (avail - h) / 2
    pad = 0.055
    fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            IN(l - pad), IN(t - pad), IN(w + 2 * pad), IN(h + 2 * pad))
    try: fr.adjustments[0] = 0.03
    except Exception: pass
    panel(fr, alpha_pct=4, line=LINE, line_w=1.0)
    s.shapes.add_picture(path, IN(l), IN(t), IN(w), IN(h))
    tf = textbox(s, l, t + h + 0.06, w, cap_h, anchor=MSO_ANCHOR.TOP)
    put(tf, [("\u25a3  " + caption, 11.5, GREY2, False, MONO)], first=True)
    for r in tf.paragraphs[0].runs:
        set_spc(r, 2)


def add_movie(s, box_l, box_t, box_w, box_h, movie, poster, caption, aspect):
    cap_h = 0.30
    avail = box_h - cap_h
    if box_w / avail > aspect:
        h = avail; w = h * aspect
    else:
        w = box_w; h = w / aspect
    l = box_l + (box_w - w) / 2
    t = box_t + (avail - h) / 2
    pad = 0.055
    fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            IN(l - pad), IN(t - pad), IN(w + 2 * pad), IN(h + 2 * pad))
    try: fr.adjustments[0] = 0.03
    except Exception: pass
    panel(fr, alpha_pct=4, line=LINE, line_w=1.0)
    s.shapes.add_movie(movie, IN(l), IN(t), IN(w), IN(h),
                       poster_frame_image=poster, mime_type="video/mp4")
    tf = textbox(s, l, t + h + 0.06, w, cap_h, anchor=MSO_ANCHOR.TOP)
    put(tf, [("\u25b6  " + caption, 11.5, GREY2, False, MONO)], first=True)
    for r in tf.paragraphs[0].runs:
        set_spc(r, 2)


def set_fullscreen(slide):
    tm = slide._element.find(qn("p:timing"))
    if tm is None:
        return
    for cmn in tm.iter(qn("p:cMediaNode")):
        cmn.set("fullScrn", "1")


def add_thumb(s, l, t, w, h, path):
    fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    try: fr.adjustments[0] = 0.06
    except Exception: pass
    panel(fr, alpha_pct=4, line=LINE, line_w=1.0)
    ins = 0.045
    s.shapes.add_picture(path, IN(l + ins), IN(t + ins),
                         IN(w - 2 * ins), IN(h - 2 * ins))


def add_stat(s, l, t, w, h, num, unit, label):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    try: sh.adjustments[0] = 0.1
    except Exception: pass
    panel(sh)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = IN(0.1)
    put(tf, [(num, 40, WHITE, True, DISP), (unit, 15, GREY, False, CJK)],
        first=True, line=1.0)
    put(tf, [(label, 13, GREY2, False, CJK)], sb=4, line=1.0)


def add_chips(s, y, texts, size=13, h=0.42, gap=0.16):
    def width(t):
        cjk = sum(1 for c in t if ord(c) > 0x2E80)
        asc = len(t) - cjk
        return 0.34 + cjk * (size / 72.0 * 1.05) + asc * (size / 72.0 * 0.6)
    ws = [width(t) for t in texts]
    total = sum(ws) + gap * (len(texts) - 1)
    x = SWIN / 2 - total / 2
    for t, w in zip(texts, ws):
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
        try: pill.adjustments[0] = 0.5
        except Exception: pass
        panel(pill, alpha_pct=5, line=LINE, line_w=1.0)
        tf = pill.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = tf.margin_bottom = 0
        put(tf, [(t, size, GREY, False, CJK)], first=True)
        x += w + gap


def add_flow(s, y, steps, h=1.25):
    n = len(steps); pw = 3.15; aw = 0.55
    total = n * pw + (n - 1) * aw
    x = SWIN / 2 - total / 2
    for i, (head, sub) in enumerate(steps):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(pw), IN(h))
        try: sh.adjustments[0] = 0.1
        except Exception: pass
        panel(sh)
        tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = IN(0.16)
        put(tf, [(head, 16, WHITE, True, DISP)], first=True, line=1.05)
        put(tf, [(sub, 11.5, GREY, False, CJK)], sb=5, line=1.2)
        x += pw
        if i < n - 1:
            atf = textbox(s, x, y, aw, h)
            put(atf, [("\u2192", 26, GREY2, True, DISP)], first=True)
            x += aw


# ---- slides -------------------------------------------------------------
def s01_cover(prs):
    s = new_slide(prs)
    cx = SWIN / 2
    ring = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              IN(cx - 0.55), IN(1.02), IN(1.1), IN(1.1))
    ring.fill.background(); ring.line.color.rgb = WHITE; ring.line.width = Pt(2.25)
    shadow_off(ring)
    try: ring.adjustments[0] = 0.22
    except Exception: pass
    rt = ring.text_frame; rt.vertical_anchor = MSO_ANCHOR.MIDDLE
    put(rt, [("\u25a3", 34, WHITE, True, DISP)], first=True)

    add_title(s, [("Blockless", WHITE)], y=2.32, size=74)
    tf = textbox(s, 0, 3.75, SWIN, 0.5)
    put(tf, [("\u786c\u4ef6\u7248   LOVABLE", 19, GREY, False, MONO)], first=True)
    for r in tf.paragraphs[0].runs: set_spc(r, 6)

    add_chips(s, 4.42, ["\u8bf4\u51fa\u6765\uff0c\u5b83\u5c31\u957f\u51fa\u6765"], size=14, h=0.56)
    add_sep(s, y=5.28)
    add_lead(s, [("\u4f60\u4e0d\u9700\u8981\u61c2\u7535\u8def\u3001\u4e0d\u9700\u8981\u4f1a\u5199\u4ee3\u7801\u3001\u4e0d\u9700\u8981\u4e70\u4efb\u4f55\u96f6\u4ef6\u3002", BODY)],
             y=5.5, size=20, w=11.5, h=0.6, line=1.35)
    add_lead(s, [("\u4e00\u53e5\u8bdd\uff0c\u5c31\u80fd\u505a\u51fa\u4e00\u4e2a\u771f\u6b63\u7684\u786c\u4ef6\u5c0f\u73a9\u610f\u3002", WHITE)],
             y=6.12, size=22, w=11.5, h=0.6, line=1.2)


def s02_gap(prs):
    s = new_slide(prs)
    add_label(s, "THE REAL GAP")
    add_title(s, [("AI \u5df2\u7ecf\u80fd\u751f\u6210\u4ee3\u7801\u548c\u7535\u8def\u56fe\u4e86\u2014\u2014\u5c0f\u767d\u8fd8\u662f\u505a\u4e0d\u51fa\u6765", WHITE)],
              y=0.98, size=27)
    add_sep(s, 1.78)
    add_card(s, 1.05, 2.02, 5.55, 1.2, "AI \u5df2\u7ecf\u505a\u5230\u7684",
             "AI \u751f\u6210\u5d4c\u5165\u5f0f\u4ee3\u7801\u3001AI \u751f\u6210\u7535\u8def\u56fe\u3001\u4e8c\u7ef4\u7535\u8def\u4eff\u771f\u2014\u2014\u542c\u8d77\u6765\u95e8\u69db\u5df2\u7ecf\u6ca1\u4e86\u3002")
    add_card(s, 6.75, 2.02, 5.55, 1.2, "\u4f46\u5c0f\u767d\u8fd8\u662f\u88ab\u6321\u5728\u95e8\u5916",
             "这些工具仍然默认你看得懂代码、会连线焊接、先买好一堆实体零件。")
    add_image(s, 1.0, 3.42, 5.6, 2.55, D("AI\u751f\u6210PCB\u548c\u4ee3\u7801.png"),
              "\u7ade\u54c1 \u00b7 AI \u751f\u6210 PCB + \u4ee3\u7801")
    add_image(s, 6.73, 3.42, 5.6, 2.55, D("wokwi.png"),
              "\u7ade\u54c1 \u00b7 Wokwi \u624b\u52a8\u4eff\u771f")
    add_lead(s, [("\u771f\u6b63\u7684\u95e8\u69db\u4e0d\u662f\u300c\u751f\u6210\u300d\uff0c\u662f\u300c\u6211\u600e\u4e48\u77e5\u9053\u5b83\u80fd\u7528\uff1f\u300d\u5c0f\u767d\u5f97\u5148\u4eb2\u773c\u770b\u5230\u5b83\u8dd1\u8d77\u6765\u3002", WHITE)],
             y=6.15, size=16, w=12, h=0.5, line=1.3)


def s03_why3d(prs):
    s = new_slide(prs)
    add_label(s, "WHY 3D & ALIVE")
    add_title(s, [("\u4e8c\u7ef4\u56fe\u7eb8\u6551\u4e0d\u4e86\u5c0f\u767d\uff0c\u4ed6\u8981\u7684\u662f\u300c\u4eb2\u773c\u770b\u5230\u5b83\u52a8\u300d", WHITE)],
              y=0.98, size=28)
    add_sep(s, 1.78)
    add_lead(s, [("AI \u662f\u5927\u8111\uff0c\u5e2e\u4f60\u628a\u4e1c\u897f\u505a\u51fa\u6765\uff1b3D \u52a8\u6001\u4eff\u771f\u662f\u773c\u775b\uff0c\u8ba9\u4f60\u770b\u5230\u5b83\u771f\u7684\u88ab\u505a\u51fa\u6765\u4e86\u3002", GREY)],
             y=2.02, size=18, w=12, h=0.6)
    add_card(s, 1.4, 2.95, 5.0, 1.75, "\u4e8c\u7ef4 \u00b7 \u9759\u6001",
             "\u4e00\u5f20\u56fe\uff0c\u770b\u5f97\u61c2\u624d\u884c \u00b7 \u4e0d\u4f1a\u52a8 \u00b7 \u300c\u5e94\u8be5\u80fd\u884c\u5427\uff1f\u300d\n\u7406\u6027\u3001\u62bd\u8c61\u3001\u5de5\u7a0b\u5e08\u601d\u7ef4\u3002")
    add_card(s, 6.93, 2.95, 5.0, 1.75, "\u4e09\u7ef4 \u00b7 \u52a8\u6001",
             "\u770b\u5f97\u89c1\u7684\u5143\u4ef6 \u00b7 \u901a\u4e0a\u7535\u771f\u7684\u8dd1 \u00b7 \u706f\u4eae\u3001\u5c4f\u5e55\u663e\u793a\u3001\u9a6c\u8fbe\u8f6c\n\u611f\u6027\u3001\u5177\u4f53\u3001\u4eba\u4eba\u90fd\u61c2\u3002", head_color=WHITE)
    add_title(s, [("\u4ece\u300c\u5e94\u8be5\u80fd\u884c\u5427\uff1f\u300d\u2192 \u300c\u5367\u69fd\uff0c\u5b83\u771f\u7684\u52a8\u4e86\uff01\u300d", WHITE)],
              y=4.95, size=30)
    add_lead(s, [("\u8fd9\u4e0d\u662f\u70ab\u6280\uff0c\u662f\u60c5\u7eea\u7684\u8f6c\u6298\u70b9\u2014\u2014\u770b\u89c1\u5b83\u6d3b\u7740\uff0c\u5c0f\u767d\u624d\u6562\u4e0b\u5355\u3001\u624d\u6562\u52a8\u624b\u3002", GREY)],
             y=5.95, size=17, w=12, h=0.5)


def s04_say(prs):
    s = new_slide(prs)
    add_label(s, "JUST SAY IT")
    add_title(s, [("\u4f60\u53ea\u9700\u8981\u8bf4\uff1a\u6211\u60f3\u8981\u4e00\u4e2a\u2026\u2026", WHITE)], y=0.98, size=30)
    add_sep(s, 1.78)
    # left: prompt box + chips ; right: image
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(0.9), IN(2.35), IN(6.3), IN(1.85))
    try: box.adjustments[0] = 0.06
    except Exception: pass
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x04, 0x08, 0x10)
    fill_alpha(box, 62); box.line.color.rgb = LINE; box.line.width = Pt(1.0); shadow_off(box)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = IN(0.3)
    put(tf, [("\u25a3 \u544a\u8bc9 Blockless \u4f60\u60f3\u505a\u4ec0\u4e48", 12, GREY2, False, MONO)],
        first=True, align=PP_ALIGN.LEFT)
    put(tf, [("\u505a\u4e00\u4e2a\u5f00\u95e8\u5c31\u4f1a\u5c16\u53eb\u7684\u76d2\u5b50\uff0c\u8fd8\u80fd\u5077\u5077\u63d0\u9192\u6211\u7684\u624b\u673a", 21, WHITE, True, CJK)],
        align=PP_ALIGN.LEFT, sb=10, line=1.35)
    # chips under prompt (left aligned-ish, reuse simple stacked line)
    ctf = textbox(s, 0.9, 4.35, 6.3, 0.8, anchor=MSO_ANCHOR.TOP)
    put(ctf, [("\u300c\u63d0\u9192\u6211\u6d47\u82b1\u7684\u5c0f\u4e1c\u897f\u300d   \u300c\u4f1a\u8ddf\u7740\u6211\u8d70\u7684\u84dd\u7259\u5c0f\u8f66\u300d", 13, GREY, False, CJK)],
        first=True, align=PP_ALIGN.LEFT, line=1.6)
    put(ctf, [("\u7528\u5927\u767d\u8bdd\uff0c\u7528\u4f60\u81ea\u5df1\u7684\u8bf4\u6cd5\u3002\u5269\u4e0b\u7684\uff0c\u4ea4\u7ed9\u5b83\u3002", 15, BODY, False, CJK)],
        align=PP_ALIGN.LEFT, sb=10)
    add_image(s, 7.5, 2.3, 5.3, 4.55, D("\u9996\u9875.png"), "\u9996\u9875")


def s05_magic(prs):
    s = new_slide(prs)
    add_label(s, "THE MAGIC MOMENT")
    add_title(s, [("\u4e0d\u662f\u7ed9\u4f60\u4e00\u5806\u4ee3\u7801\u3002\u662f\u8ba9\u4f60\u770b\u89c1\u2014\u2014\u5b83\uff0c\u52a8\uff0c\u4e86\u3002", WHITE)],
              y=0.9, size=26)
    add_movie(s, 0.7, 1.82, 11.93, 5.32, MOVIE, D("\u5b9e\u9645\u6f14\u793a.png"),
              "3D \u4eff\u771f\u6f14\u793a \u00b7 \u70b9\u51fb\u5168\u5c4f\u64ad\u653e", 1.6)
    set_fullscreen(s)


def s06_engine(prs):
    s = new_slide(prs)
    add_label(s, "THE 3D ENGINE")
    add_title(s, [("\u4e0d\u662f\u52a8\u753b\uff0c\u662f\u4eff\u771f", WHITE)], y=0.98, size=32)
    add_sep(s, 1.78)
    # left column text+cards, right image
    add_lead(s, [("\u4f60\u7684\u786c\u4ef6\u4ee5\u771f\u5b9e\u7684 3D \u51fa\u73b0\u5728\u6d4f\u89c8\u5668\u91cc\uff0c\u53ef\u4ee5\u8f6c\u3001\u53ef\u4ee5\u62c6\uff0c\u901a\u4e0a\u7535\u5c31\u771f\u7684\u8dd1\u8d77\u6765\u3002", GREY)],
             y=2.1, size=16, w=6.4, h=0.9)
    # override lead centering by placing manually left
    add_card(s, 0.9, 3.1, 5.9, 1.4, "\u770b\u5f97\u89c1",
             "\u706f\u4eae\u3001\u5c4f\u5e55\u663e\u793a\u3001\u9a6c\u8fbe\u8f6c\u3001\u4f20\u611f\u5668\u8bfb\u6570\u3001\u4e32\u53e3\u6eda\u52a8\u2014\u2014\u90fd\u5728 3D \u91cc\u5b9e\u65f6\u53d1\u751f\u3002")
    add_card(s, 0.9, 4.65, 5.9, 1.4, "\u662f\u771f\u7684",
             "\u8dd1\u7684\u662f\u540c\u4e00\u4efd\u56fa\u4ef6\uff1a\u865a\u62df\u786c\u4ef6\u548c\u771f\u786c\u4ef6\u8868\u73b0\u4e00\u81f4\uff0c\u4e0d\u662f\u9884\u5f55\u52a8\u753b\u3002", head_color=WHITE)
    add_image(s, 7.0, 2.15, 5.8, 4.6, D("\u4eff\u771f\u6548\u679c.png"), "\u4eff\u771f\u6548\u679c")


def s07_first(prs):
    s = new_slide(prs)
    add_label(s, "WORLD FIRST")
    add_title(s, [("\u4e16\u754c\u4e0a\u7b2c\u4e00\u4e2a\u300c\u4e00\u53e5\u8bdd \u2192 \u4e00\u6574\u4e2a 3D \u786c\u4ef6 \u2192 \u4e00\u952e\u53d8\u5b9e\u7269\u300d", WHITE)],
              y=0.98, size=27)
    add_sep(s, 1.78)
    top = [("AI \u753b\u7535\u8def\u56fe", "\u4e00\u5f20\u4e8c\u7ef4\u56fe\uff0c\u8fd8\u5f97\u4f60\u770b\u61c2\u3001\u8fde\u7ebf\u3001\u710a\u63a5"),
           ("\u624b\u52a8\u4eff\u771f \u00b7 Wokwi", "\u5f88\u771f\uff0c\u4f46\u8981\u4f60\u62d6\u96f6\u4ef6\u3001\u5199\u4ee3\u7801\u3001\u8c03\u8bd5"),
           ("AI \u5199\u5d4c\u5165\u5f0f\u4ee3\u7801", "\u7ed9\u4f60\u4ee3\u7801\uff0c\u80fd\u4e0d\u80fd\u8dd1\u81ea\u5df1\u60f3\u529e\u6cd5")]
    bot = [("\u8bf4\u4e00\u53e5\u8bdd\u5c31\u751f\u6210", "\u800c\u4e0d\u662f\u624b\u52a8\u62d6\u96f6\u4ef6\u3001\u5199\u4ee3\u7801"),
           ("\u770b\u5f97\u89c1\u5730 3D \u8dd1\u8d77\u6765", "\u52a8\u6001\u4eff\u771f\uff0c\u4e0d\u662f\u9759\u6001\u6e32\u67d3\u56fe"),
           ("\u4e00\u952e\u53d8\u6210\u771f\u4e1c\u897f", "\u540c\u4e00\u4efd\u8bbe\u8ba1\uff0c\u76f4\u63a5\u4e0b\u5355\u9020\u51fa\u6765")]
    w = 3.75; gap = 0.2; total = 3 * w + 2 * gap; x0 = SWIN / 2 - total / 2
    for i, (h, b) in enumerate(top):
        add_card(s, x0 + i * (w + gap), 2.1, w, 1.5, h, b, head_color=GREY, head_size=15)
    for i, (h, b) in enumerate(bot):
        add_card(s, x0 + i * (w + gap), 3.78, w, 1.5, h, b, head_color=WHITE, head_size=15)
    add_lead(s, [("Lovable \u8ba9\u4e0d\u61c2\u4ee3\u7801\u7684\u4eba\u505a\u51fa\u8f6f\u4ef6\u3002Blockless \u8ba9\u4e0d\u61c2\u786c\u4ef6\u7684\u4eba\u505a\u51fa\u786c\u4ef6\u3002", WHITE)],
             y=5.6, size=18, w=12, h=0.6)


def s08_cases(prs):
    s = new_slide(prs)
    add_label(s, "USER CASES")
    add_title(s, [("\u771f\u5b9e\u7528\u6237\uff0c\u771f\u5b9e\u505a\u51fa\u6765\u7684\u4e1c\u897f", WHITE)], y=0.96, size=28)
    add_lead(s, [("\u4e00\u53e5\u8bdd\u8bf4\u51fa\u60f3\u6cd5\uff0c\u4e5d\u4e2a\u4eba\u505a\u51fa\u4e86\u4e5d\u4e2a\u771f\u5bb6\u4f19\u2014\u2014\u4ece\u4eff\u771f\u5230\u5b9e\u7269\uff0c\u5168\u5728\u8fd9\u3002", GREY)],
             y=1.72, size=15, w=12, h=0.45)
    cw, ch, g = 2.05, 1.54, 0.16
    gw = 3 * cw + 2 * g
    x0 = SWIN / 2 - gw / 2
    y0 = 2.35
    for i in range(9):
        r, c = divmod(i, 3)
        add_thumb(s, x0 + c * (cw + g), y0 + r * (ch + g), cw, ch,
                  os.path.join(THUMB, "case%d.jpg" % (i + 1)))


def s09_rednote(prs):
    s = new_slide(prs)
    add_label(s, "COLD START \u00b7 \u5c0f\u7ea2\u4e66")
    add_title(s, [("\u521a\u505a\u5b8c\uff0c\u81ea\u5df1\u968f\u624b\u53d1\u4e86 6 \u6761\uff0c\u6570\u636e\u5c31\u6765\u4e86", WHITE)], y=0.98, size=28)
    add_sep(s, 1.78)
    add_lead(s, [("\u4ea7\u54c1\u521a\u5f00\u53d1\u5b8c\uff0c\u6211\u4eec\u81ea\u5df1\u5728\u5c0f\u7ea2\u4e66\u53d1\u4e86 6 \u6761\u7b14\u8bb0\uff0c\u4e00\u5206\u94b1\u6ca1\u63a8\u3002", GREY)],
             y=2.02, size=16, w=12, h=0.45)
    # left 2x2 stats
    sx, sy, sw, sh, gg = 0.95, 2.75, 2.75, 1.55, 0.25
    add_stat(s, sx, sy, sw, sh, "1500", "\u6b21", "\u603b\u6d4f\u89c8\u91cf")
    add_stat(s, sx + sw + gg, sy, sw, sh, "60", "\u4e2a", "\u603b\u70b9\u8d5e\u6570")
    add_stat(s, sx, sy + sh + gg, sw, sh, "81", "\u4e2a", "\u603b\u6536\u85cf\u6570")
    add_stat(s, sx + sw + gg, sy + sh + gg, sw, sh, "9", "\u6761", "\u603b\u8bc4\u8bba\u6570")
    add_image(s, 7.1, 2.55, 5.6, 4.3, D("\u5c0f\u7ea2\u4e66\u6d4f\u89c8\u91cf.png"), "\u5c0f\u7ea2\u4e66\u6d4f\u89c8\u91cf")
    add_lead(s, [("\u6536\u85cf(81) \u751a\u81f3\u9ad8\u8fc7\u70b9\u8d5e(60)\u2014\u2014\u770b\u7684\u4eba\u662f\u771f\u60f3\u7559\u7740\u81ea\u5df1\u505a\u3002", WHITE)],
             y=6.5, size=15, w=8, h=0.4)


def s10_end(prs):
    s = new_slide(prs)
    cx = SWIN / 2
    ring = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(cx - 0.45), IN(1.15), IN(0.9), IN(0.9))
    ring.fill.background(); ring.line.color.rgb = WHITE; ring.line.width = Pt(2.0); shadow_off(ring)
    try: ring.adjustments[0] = 0.22
    except Exception: pass
    rt = ring.text_frame; rt.vertical_anchor = MSO_ANCHOR.MIDDLE
    put(rt, [("\u25a3", 30, WHITE, True, DISP)], first=True)
    add_title(s, [("\u63cf\u8ff0\u4f60\u7684\u60f3\u6cd5 \u2192 \u770b\u5b83\u5728\u6d4f\u89c8\u5668\u91cc\u8dd1 \u2192 \u4e00\u952e\u9020\u51fa\u771f\u4e1c\u897f\u3002", WHITE)],
              y=2.5, size=30)
    add_lead(s, [("\u300cAI \u5199\u4ee3\u7801\u300d\u662f\u53bb\u5e74\u7684\u6545\u4e8b\u3002\u4eca\u5e74\u7684\u6545\u4e8b\u662f\uff1a\u8bf4\u4e00\u53e5\u8bdd\uff0c\u4e00\u4e2a\u786c\u4ef6\u5728\u4f60\u773c\u524d\u6d3b\u8fc7\u6765\u3002", GREY)],
             y=3.5, size=18, w=11.5, h=0.9)
    add_lead(s, [("Blockless\uff0c\u662f\u8fd9\u4e2a\u6545\u4e8b\u91cc\u552f\u4e00\u7684\u73a9\u5bb6\u3002", WHITE)],
             y=4.35, size=20, w=11.5, h=0.6)
    add_chips(s, 5.35, ["\u786c\u4ef6\u7248 Lovable", "Vibe Hardware", "\u8bf4\u51fa\u6765\uff0c\u5b83\u5c31\u957f\u51fa\u6765"],
              size=14, h=0.5)


def main():
    make_bg(); make_thumbs()
    prs = Presentation()
    prs.slide_width = IN(SWIN); prs.slide_height = IN(SHIN)
    for fn in (s01_cover, s02_gap, s03_why3d, s04_say, s05_magic,
               s06_engine, s07_first, s08_cases, s09_rednote, s10_end):
        fn(prs)
    out = OUT
    try:
        prs.save(out)
    except PermissionError:
        out = os.path.join(ROOT, "Blockless-new.pptx")
        prs.save(out)
        print("NOTE: Blockless.pptx was locked (open in PowerPoint?)")
    print("saved", out, "slides", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
